# Script para arquivos numerados (0 a 10)
# Salve como criar_slides_numerados.py

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
import re

print("=" * 60)
print("📊 CRIAÇÃO DE APRESENTAÇÃO COM IMAGENS NUMERADAS")
print("=" * 60)

# Configurações
PASTA_CAPTURAS = r"C:\Users\IFAM-CBDA\Pictures\Screenshots\python2presentation"
NOME_ARQUIVO = "marp_Presentation.pptx"
REPO_URL = "https://github.com/seu-usuario/seu-repositorio"  # Altere para seu repo

def extrair_numero(nome_arquivo):
    """Extrai o número do início do nome do arquivo"""
    match = re.search(r'^(\d+)', nome_arquivo)
    return int(match.group(1)) if match else 999

def nome_para_titulo(nome_arquivo):
    """Converte '01_readme_example' para 'Readme Example'"""
    nome_sem_ext = os.path.splitext(nome_arquivo)[0]
    # Remove o número no início
    nome_sem_numero = re.sub(r'^\d+_', '', nome_sem_ext)
    # Substitui underscores por espaços e capitaliza
    return nome_sem_numero.replace('_', ' ').title()

def criar_apresentacao():
    # Verifica se a pasta existe
    if not os.path.exists(PASTA_CAPTURAS):
        print(f"❌ ERRO: Pasta '{PASTA_CAPTURAS}' não encontrada!")
        print(f"Crie uma pasta chamada '{PASTA_CAPTURAS}' na mesma pasta deste script.")
        print(f"Coloque suas imagens numeradas dentro (0.png, 1.jpg, etc.)")
        input("\nPressione Enter para sair...")
        return
    
    # Lista e ordena arquivos pelo número no início
    imagens = []
    for arquivo in os.listdir(PASTA_CAPTURAS):
        if arquivo.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp')):
            imagens.append(arquivo)
    
    if len(imagens) == 0:
        print("❌ Nenhuma imagem encontrada na pasta 'capturas'!")
        print("Certifique-se de que os arquivos estão com extensões:")
        print(".png, .jpg, .jpeg, .gif ou .bmp")
        input("\nPressione Enter para sair...")
        return
    
    # Ordena pelo número extraído do nome
    imagens.sort(key=extrair_numero)
    
    print(f"\n✅ Encontradas {len(imagens)} imagens:")
    for i, img in enumerate(imagens, 1):
        print(f"  {i:2d}. {img}")
    
    # Cria a apresentação
    prs = Presentation()
    
    # Tamanho widescreen (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ========== SLIDE 1: CAPA ==========
    print("\n📄 Criando slide de capa...")
    slide_capa = prs.slides.add_slide(prs.slide_layouts[0])
    slide_capa.shapes.title.text = "GitHub Repository Deep Dive"
    slide_capa.placeholders[1].text = f"Repository: {REPO_URL}\nFocus: Learning Technical English"
    
    # ========== SLIDE 2: AGENDA ==========
    print("📄 Criando slide de agenda...")
    slide_agenda = prs.slides.add_slide(prs.slide_layouts[1])
    slide_agenda.shapes.title.text = "Today's Agenda"
    conteudo_agenda = slide_agenda.placeholders[1]
    
    # Cria agenda baseada nos nomes dos arquivos
    texto_agenda = ""
    for i, img in enumerate(imagens, 1):
        titulo = nome_para_titulo(img)
        texto_agenda += f"{i}. {titulo}\n"
    
    conteudo_agenda.text = texto_agenda
    
    # ========== SLIDES COM IMAGENS ==========
    print("\n🖼️ Adicionando slides com imagens...")
    
    for i, nome_imagem in enumerate(imagens, 1):
        print(f"  Slide {i+2}: {nome_imagem}")
        
        # Cria slide em branco
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Título do slide (baseado no nome do arquivo)
        titulo_slide = nome_para_titulo(nome_imagem)
        caixa_titulo = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), 
                                                Inches(12), Inches(0.7))
        frame_titulo = caixa_titulo.text_frame
        frame_titulo.text = f"Step {i}: {titulo_slide}"
        frame_titulo.paragraphs[0].font.size = Pt(28)
        frame_titulo.paragraphs[0].font.bold = True
        
        # Adiciona a imagem
        caminho_completo = os.path.join(PASTA_CAPTURAS, nome_imagem)
        
        try:
            # Calcula posição para centralizar
            img = slide.shapes.add_picture(
                caminho_completo,
                Inches(1),      # Margem esquerda
                Inches(1.2),    # Topo abaixo do título
                width=Inches(11.33)  # Largura fixa
            )
            
            # ========== LEGENDA EM INGLÊS ==========
            # Adiciona caixa de vocabulário técnico
            caixa_vocabulario = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5), Inches(12), Inches(0.8)
            )
            frame_vocab = caixa_vocabulario.text_frame
            
            # Diferentes dicas baseadas no número do slide
            dicas_ingles = [
                "🔤 Key terms: repository, readme, documentation, license",
                "💻 Code vocabulary: function, variable, loop, condition, class",
                "🔄 Git commands: clone, commit, push, pull, branch, merge",
                "🎨 UI terms: interface, button, menu, input, output, display",
                "🐛 Debugging: issue, bug, error, warning, log, trace",
                "📈 Project: milestone, release, version, update, patch",
                "👥 Collaboration: contributor, review, comment, suggestion",
                "⚙️ Configuration: settings, options, preferences, setup",
                "📁 Structure: folder, directory, file, path, extension",
                "🔍 Navigation: search, filter, sort, find, browse"
            ]
            
            idx = min(i-1, len(dicas_ingles)-1)
            frame_vocab.text = "📚 English Practice: " + dicas_ingles[idx]
            frame_vocab.paragraphs[0].font.size = Pt(14)
            frame_vocab.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
            
        except Exception as e:
            print(f"  ⚠️ Erro ao adicionar {nome_imagem}: {str(e)}")
    
    # ========== SLIDE FINAL: PRÁTICA ==========
    print("\n📄 Criando slide final de prática...")
    slide_final = prs.slides.add_slide(prs.slide_layouts[5])
    slide_final.shapes.title.text = "Your English Learning Challenge"
    
    conteudo_final = slide_final.placeholders[1]
    texto_final = f"""Your Mission:

1. 📖 Read the README.md at: {REPO_URL}
2. 🔍 Identify 10 technical words you didn't know
3. ✍️ Write a simple issue in English
4. 🗣️ Explain this project to someone in English

Remember:
• Technical English uses simple grammar
• Many words are similar to Portuguese
• Practice daily for 15 minutes
• Don't be afraid to make mistakes!

🚀 Start your open-source journey today!"""
    
    conteudo_final.text = texto_final
    
    # ========== SALVAR ==========
    prs.save(NOME_ARQUIVO)
    
    # ========== RESUMO ==========
    print("\n" + "=" * 60)
    print("🎉 APRESENTAÇÃO CRIADA COM SUCESSO!")
    print("=" * 60)
    print(f"📂 Arquivo gerado: {NOME_ARQUIVO}")
    print(f"📊 Total de slides: {len(prs.slides)}")
    print(f"🖼️ Imagens incluídas: {len(imagens)}")
    print(f"📍 Slides com imagens: {len(imagens)}")
    print(f"📍 Slides adicionais: 3 (capa, agenda, prática)")
    print("\n📌 Estrutura da apresentação:")
    print("  1. Capa")
    print("  2. Agenda")
    for i, img in enumerate(imagens, 3):
        print(f"  {i}. {nome_para_titulo(img)}")
    print(f"  {len(imagens)+3}. Prática de Inglês")
    print("\n✅ Pronto! Abra o arquivo no PowerPoint para personalizar.")
    print("=" * 60)
    
    # Pergunta se quer abrir
    abrir = input("\nAbrir a apresentação agora? (s/n): ").lower()
    if abrir == 's':
        os.startfile(NOME_ARQUIVO)
        print("Abrindo PowerPoint...")

if __name__ == "__main__":
    criar_apresentacao()
    input("\nPressione Enter para sair...")